from __future__ import annotations

import json
import re
import uuid
from collections.abc import Iterator
from pathlib import Path
from typing import Any

from app.detectors import JapanesePiiEngine, apply_mask
from app.models import DictionaryEntry, DocumentBlock, Finding


SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".pdf"}
MAX_UPLOAD_BYTES = 25 * 1024 * 1024


class UnsupportedDocumentError(ValueError):
    pass


class DocumentSessionStore:
    def __init__(self, root: Path) -> None:
        self.root = root
        self.root.mkdir(parents=True, exist_ok=True)

    def create(self, filename: str, content: bytes) -> tuple[str, Path]:
        extension = Path(filename).suffix.lower()
        if extension not in SUPPORTED_EXTENSIONS:
            raise UnsupportedDocumentError("DOCX、PPTX、PDFのみアップロードできます。")
        if len(content) > MAX_UPLOAD_BYTES:
            raise UnsupportedDocumentError("ファイルサイズは25MB以下にしてください。")
        if not content:
            raise UnsupportedDocumentError("空のファイルは処理できません。")
        session_id = uuid.uuid4().hex
        session_dir = self.root / session_id
        session_dir.mkdir(mode=0o700)
        safe_name = f"input{extension}"
        source = session_dir / safe_name
        source.write_bytes(content)
        metadata = {
            "session_id": session_id,
            "original_filename": _safe_download_name(filename),
            "source": safe_name,
            "extension": extension,
        }
        (session_dir / "session.json").write_text(
            json.dumps(metadata, ensure_ascii=False), encoding="utf-8"
        )
        return session_id, source

    def save_analysis(self, session_id: str, blocks: list[DocumentBlock], findings: list[Finding]) -> None:
        payload = {
            "blocks": [block.model_dump() for block in blocks],
            "findings": [finding.model_dump() for finding in findings],
        }
        (self._session_dir(session_id) / "analysis.json").write_text(
            json.dumps(payload, ensure_ascii=False), encoding="utf-8"
        )

    def load(self, session_id: str) -> tuple[dict[str, Any], list[DocumentBlock], list[Finding]]:
        session_dir = self._session_dir(session_id)
        metadata = json.loads((session_dir / "session.json").read_text(encoding="utf-8"))
        analysis = json.loads((session_dir / "analysis.json").read_text(encoding="utf-8"))
        blocks = [DocumentBlock(**item) for item in analysis["blocks"]]
        findings = [Finding(**item) for item in analysis["findings"]]
        return metadata, blocks, findings

    def source_path(self, session_id: str, metadata: dict[str, Any]) -> Path:
        return self._session_dir(session_id) / metadata["source"]

    def output_path(self, session_id: str, extension: str) -> Path:
        return self._session_dir(session_id) / f"masked{extension}"

    def _session_dir(self, session_id: str) -> Path:
        if not re.fullmatch(r"[0-9a-f]{32}", session_id):
            raise FileNotFoundError(session_id)
        path = self.root / session_id
        if not path.is_dir():
            raise FileNotFoundError(session_id)
        return path


class DocumentProcessor:
    def extract(self, path: Path) -> list[DocumentBlock]:
        extension = path.suffix.lower()
        if extension == ".docx":
            return self._extract_docx(path)
        if extension == ".pptx":
            return self._extract_pptx(path)
        if extension == ".pdf":
            return self._extract_pdf(path)
        raise UnsupportedDocumentError("未対応のファイル形式です。")

    def mask(
        self,
        source: Path,
        output: Path,
        findings: list[Finding],
        accepted_ids: set[str],
        mask_character: str,
    ) -> None:
        by_block: dict[str, list[Finding]] = {}
        for finding in findings:
            if finding.id in accepted_ids:
                by_block.setdefault(finding.block_id, []).append(finding)
        extension = source.suffix.lower()
        if extension == ".docx":
            self._mask_docx(source, output, by_block, accepted_ids, mask_character)
        elif extension == ".pptx":
            self._mask_pptx(source, output, by_block, accepted_ids, mask_character)
        elif extension == ".pdf":
            self._mask_pdf(source, output, by_block, accepted_ids)
        else:
            raise UnsupportedDocumentError("未対応のファイル形式です。")

    @staticmethod
    def analyze_blocks(
        blocks: list[DocumentBlock],
        engine: JapanesePiiEngine,
        entities: list[str] | None,
        dictionary: list[DictionaryEntry],
    ) -> list[Finding]:
        findings: list[Finding] = []
        for block in blocks:
            findings.extend(engine.analyze(block.text, entities, dictionary, block.id))
        return findings

    def _extract_docx(self, path: Path) -> list[DocumentBlock]:
        document = _load_docx(path)
        return [
            DocumentBlock(id=block_id, text=paragraph.text)
            for block_id, paragraph in _iter_docx_paragraphs(document)
            if paragraph.text.strip()
        ]

    def _mask_docx(
        self,
        source: Path,
        output: Path,
        by_block: dict[str, list[Finding]],
        accepted_ids: set[str],
        mask_character: str,
    ) -> None:
        document = _load_docx(source)
        for block_id, paragraph in _iter_docx_paragraphs(document):
            if block_id in by_block:
                _mask_text_runs(
                    paragraph.runs,
                    paragraph.text,
                    by_block[block_id],
                    accepted_ids,
                    mask_character,
                )
        document.save(output)

    def _extract_pptx(self, path: Path) -> list[DocumentBlock]:
        presentation = _load_pptx(path)
        return [
            DocumentBlock(id=block_id, text=text_frame.text)
            for block_id, text_frame in _iter_pptx_text_frames(presentation)
            if text_frame.text.strip()
        ]

    def _mask_pptx(
        self,
        source: Path,
        output: Path,
        by_block: dict[str, list[Finding]],
        accepted_ids: set[str],
        mask_character: str,
    ) -> None:
        presentation = _load_pptx(source)
        for block_id, text_frame in _iter_pptx_text_frames(presentation):
            findings = by_block.get(block_id)
            if not findings:
                continue
            full_text = text_frame.text
            masked = apply_mask(full_text, findings, accepted_ids, mask_character)
            offset = 0
            for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
                for run in paragraph.runs:
                    length = len(run.text)
                    run.text = masked[offset : offset + length]
                    offset += length
                if paragraph_index < len(text_frame.paragraphs) - 1:
                    offset += 1  # text_frame.text joins paragraphs with a newline
        presentation.save(output)

    def _extract_pdf(self, path: Path) -> list[DocumentBlock]:
        reader, ContentStream, TextStringObject = _load_pdf_types(path)
        blocks: list[DocumentBlock] = []
        for page_index, page in enumerate(reader.pages):
            content = page.get_contents()
            if content is None:
                continue
            stream = ContentStream(content, reader)
            for operation_index, (operands, operator) in enumerate(stream.operations):
                if operator == b"Tj" and operands and isinstance(operands[0], TextStringObject):
                    text = str(operands[0])
                    if text.strip():
                        blocks.append(DocumentBlock(id=f"pdf:{page_index}:{operation_index}:0", text=text))
                elif operator == b"TJ" and operands:
                    for item_index, item in enumerate(operands[0]):
                        if isinstance(item, TextStringObject) and str(item).strip():
                            blocks.append(
                                DocumentBlock(
                                    id=f"pdf:{page_index}:{operation_index}:{item_index}",
                                    text=str(item),
                                )
                            )
        return blocks

    def _mask_pdf(
        self,
        source: Path,
        output: Path,
        by_block: dict[str, list[Finding]],
        accepted_ids: set[str],
    ) -> None:
        from pypdf import PdfReader, PdfWriter
        from pypdf.generic import ContentStream, TextStringObject

        reader = PdfReader(source)
        writer = PdfWriter()
        for page_index, page in enumerate(reader.pages):
            content = page.get_contents()
            if content is not None:
                stream = ContentStream(content, reader)
                for operation_index, (operands, operator) in enumerate(stream.operations):
                    if operator == b"Tj" and operands and isinstance(operands[0], TextStringObject):
                        block_id = f"pdf:{page_index}:{operation_index}:0"
                        if block_id in by_block:
                            original = str(operands[0])
                            operands[0] = TextStringObject(
                                apply_mask(original, by_block[block_id], accepted_ids, "*")
                            )
                    elif operator == b"TJ" and operands:
                        for item_index, item in enumerate(operands[0]):
                            block_id = f"pdf:{page_index}:{operation_index}:{item_index}"
                            if block_id in by_block and isinstance(item, TextStringObject):
                                operands[0][item_index] = TextStringObject(
                                    apply_mask(str(item), by_block[block_id], accepted_ids, "*")
                                )
                page.replace_contents(stream)
            writer.add_page(page)
        with output.open("wb") as file_handle:
            writer.write(file_handle)


def _load_docx(path: Path):
    from docx import Document

    return Document(path)


def _iter_docx_paragraphs(document) -> Iterator[tuple[str, Any]]:
    seen_parts: set[str] = set()

    def walk(container, prefix: str) -> Iterator[tuple[str, Any]]:
        paragraph_index = 0
        table_index = 0
        for item in container.iter_inner_content():
            if item.__class__.__name__ == "Paragraph":
                yield f"{prefix}:p:{paragraph_index}", item
                paragraph_index += 1
            else:
                from docx.table import _Cell

                current_table = table_index
                table_index += 1
                for row_index, row in enumerate(item.rows):
                    for cell_index, table_cell in enumerate(row._tr.tc_lst):
                        cell = _Cell(table_cell, item)
                        yield from walk(
                            cell,
                            f"{prefix}:t:{current_table}:r:{row_index}:c:{cell_index}",
                        )

    yield from walk(document, "doc")
    for section_index, section in enumerate(document.sections):
        containers = [
            ("header", section.header),
            ("first-header", section.first_page_header),
            ("even-header", section.even_page_header),
            ("footer", section.footer),
            ("first-footer", section.first_page_footer),
            ("even-footer", section.even_page_footer),
        ]
        for name, container in containers:
            part_name = str(container.part.partname)
            if part_name in seen_parts:
                continue
            seen_parts.add(part_name)
            yield from walk(container, f"section:{section_index}:{name}")


def _mask_text_runs(
    runs: list[Any],
    full_text: str,
    findings: list[Finding],
    accepted_ids: set[str],
    mask_character: str,
) -> None:
    masked = apply_mask(full_text, findings, accepted_ids, mask_character)
    offset = 0
    for run in runs:
        length = len(run.text)
        run.text = masked[offset : offset + length]
        offset += length


def _load_pptx(path: Path):
    from pptx import Presentation

    return Presentation(path)


def _iter_pptx_text_frames(presentation) -> Iterator[tuple[str, Any]]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk_shapes(shapes, prefix: str) -> Iterator[tuple[str, Any]]:
        for shape_index, shape in enumerate(shapes):
            shape_prefix = f"{prefix}:shape:{shape_index}"
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk_shapes(shape.shapes, shape_prefix)
            if getattr(shape, "has_text_frame", False):
                yield f"{shape_prefix}:text", shape.text_frame
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows):
                    for cell_index, cell in enumerate(row.cells):
                        yield f"{shape_prefix}:r:{row_index}:c:{cell_index}", cell.text_frame

    for slide_index, slide in enumerate(presentation.slides):
        yield from walk_shapes(slide.shapes, f"slide:{slide_index}")
        if slide.has_notes_slide:
            yield from walk_shapes(slide.notes_slide.shapes, f"slide:{slide_index}:notes")


def _load_pdf_types(path: Path):
    from pypdf import PdfReader
    from pypdf.generic import ContentStream, TextStringObject

    return PdfReader(path), ContentStream, TextStringObject


def _safe_download_name(filename: str) -> str:
    name = Path(filename).name
    cleaned = re.sub(r"[^\w.()\-\u3000-\u30ff\u3400-\u9fff]", "_", name)
    return cleaned[:180] or "document"
