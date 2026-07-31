from app.detectors import JapanesePiiEngine
from app.documents import DocumentProcessor
from app.models import DictionaryEntry


def entry(term: str) -> DictionaryEntry:
    return DictionaryEntry(
        id=1,
        term=term,
        entity_type="CUSTOM",
        note="",
        created_at="2026-01-01T00:00:00+00:00",
    )


def test_docx_masks_a_term_split_across_runs_and_keeps_other_text(tmp_path) -> None:
    from docx import Document

    source = tmp_path / "source.docx"
    output = tmp_path / "masked.docx"
    document = Document()
    paragraph = document.add_paragraph()
    paragraph.add_run("担当は山田")
    paragraph.add_run("太郎です")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "連絡先 taro@example.jp"
    document.sections[0].header.paragraphs[0].text = "社外秘"
    document.save(source)

    processor = DocumentProcessor()
    blocks = processor.extract(source)
    engine = JapanesePiiEngine()
    findings = processor.analyze_blocks(
        blocks,
        engine,
        ["CUSTOM", "EMAIL_ADDRESS"],
        [entry("山田太郎"), entry("社外秘")],
    )
    processor.mask(source, output, findings, {item.id for item in findings}, "█")

    masked = Document(output)
    assert masked.paragraphs[0].text == "担当は████です"
    assert masked.tables[0].cell(0, 0).text == "連絡先 ███████████████"
    assert masked.sections[0].header.paragraphs[0].text == "███"


def test_pptx_masks_a_term_split_across_runs_and_table_text(tmp_path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    source = tmp_path / "source.pptx"
    output = tmp_path / "masked.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.add_run().text = "顧客：山田"
    paragraph.add_run().text = "太郎"
    table_shape = slide.shapes.add_table(1, 1, Inches(1), Inches(2), Inches(5), Inches(1))
    table_shape.table.cell(0, 0).text = "taro@example.jp"
    presentation.save(source)

    processor = DocumentProcessor()
    blocks = processor.extract(source)
    findings = processor.analyze_blocks(
        blocks,
        JapanesePiiEngine(),
        ["CUSTOM", "EMAIL_ADDRESS"],
        [entry("山田太郎")],
    )
    processor.mask(source, output, findings, {item.id for item in findings}, "█")

    masked = Presentation(output)
    assert masked.slides[0].shapes[0].text == "顧客：████"
    assert masked.slides[0].shapes[1].table.cell(0, 0).text == "███████████████"


def test_pdf_replaces_content_stream_text_instead_of_overlaying_it(tmp_path) -> None:
    from pypdf import PdfReader, PdfWriter
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    source = tmp_path / "source.pdf"
    output = tmp_path / "masked.pdf"
    writer = PdfWriter()
    page = writer.add_blank_page(width=300, height=300)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_reference = writer._add_object(font)
    page[NameObject("/Resources")] = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_reference})}
    )
    stream = DecodedStreamObject()
    stream.set_data(b"BT /F1 12 Tf 20 150 Td (taro@example.jp) Tj ET")
    page[NameObject("/Contents")] = writer._add_object(stream)
    with source.open("wb") as file_handle:
        writer.write(file_handle)

    processor = DocumentProcessor()
    blocks = processor.extract(source)
    findings = processor.analyze_blocks(
        blocks, JapanesePiiEngine(), ["EMAIL_ADDRESS"], []
    )
    processor.mask(source, output, findings, {item.id for item in findings}, "█")

    extracted = PdfReader(output).pages[0].extract_text()
    assert "taro@example.jp" not in extracted
    assert "***************" in extracted
    assert b"taro@example.jp" not in output.read_bytes()

